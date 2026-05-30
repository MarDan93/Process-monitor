import dash
import dash_bootstrap_components as dbc
from dash import dcc, html, Input, Output, State, callback, no_update
import os

from core.ai_client import call_ai, detect_language

app = dash.Dash(
    __name__,
    use_pages=True,
    external_stylesheets=[dbc.themes.BOOTSTRAP],
    suppress_callback_exceptions=True,
    title="Process Monitor",
    update_title=None,
)
server = app.server  # for gunicorn

# ── Load API keys from environment (fallback for local dev) ──
for _var in ("GEMINI_API_KEY", "ANTHROPIC_API_KEY"):
    if not os.environ.get(_var):
        try:
            from dotenv import load_dotenv
            load_dotenv()
        except ImportError:
            pass
        break

# ── Sidebar nav items ─────────────────────────────────────────
NAV_ITEMS = [
    ("/",           "🏠", "Dashboard"),
    ("/setup",      "🚀", "Setup"),
    ("/training",   "📐", "Training"),
    ("/structure",  "📊", "Structure"),
    ("/monitoring", "🔍", "Monitoring"),
    ("/analysis",   "📋", "Root Cause"),
]

PAGE_CHIPS = {
    "/":           ["Qual è lo stato attuale?", "Cosa devo fare ora?"],
    "/setup":      ["Quali variabili sono più importanti?", "Come scelgo le Y?"],
    "/training":   ["Quante componenti scegliere?", "Il modello è ben calibrato?"],
    "/structure":  ["Cosa rappresentano PC1 e PC2?", "Quali variabili correlano?"],
    "/monitoring": ["Perché questi cicli sono anomali?", "Quali variabili controllare?"],
    "/analysis":   ["Qual è la root cause principale?", "Cosa devo monitorare?"],
}


def _sidebar():
    return html.Div(
        id="sidebar",
        children=[
            html.Div("Navigazione", className="nav-section"),
            *[
                dcc.Link(
                    [
                        html.Span(icon, className="nav-icon"),
                        html.Span(label),
                        html.Span(id={"type": "nav-status", "page": path}, className="nav-status"),
                    ],
                    href=path,
                    className="nav-link",
                    id={"type": "nav-link", "page": path},
                )
                for path, icon, label in NAV_ITEMS
            ],
            html.Hr(style={"margin": "12px 16px", "borderColor": "#E2E8F0"}),
            html.Div("Modello", className="nav-section"),
            dbc.Button(
                "⬇ Salva sessione",
                id="btn-save-session",
                color="light",
                size="sm",
                className="mx-3 mb-2",
                style={"width": "calc(100% - 24px)", "fontSize": "12px"},
            ),
            dcc.Download(id="download-session"),
            dcc.Upload(
                id="upload-session",
                children=dbc.Button(
                    "⬆ Carica sessione",
                    color="light",
                    size="sm",
                    style={"width": "100%", "fontSize": "12px"},
                ),
                style={"padding": "0 12px"},
                multiple=False,
            ),
            html.Div(id="session-load-msg", className="px-3 mt-1"),
        ],
    )


def _header():
    return html.Div(
        id="app-header",
        children=[
            html.Span("🏭", style={"fontSize": "20px"}),
            html.Span(
                "Process Monitor",
                style={"fontWeight": 700, "fontSize": "17px", "color": "#1E293B"},
            ),
            html.Span(
                id="header-process-name",
                style={"fontSize": "12px", "color": "#64748B", "marginLeft": "4px"},
            ),
            html.Div(className="header-spacer"),
            html.Div(id="header-status-badge"),
            dbc.Button(
                "⬇ Esporta",
                id="btn-export-menu",
                color="light",
                size="sm",
                style={"fontSize": "12px", "marginRight": "6px"},
            ),
            dbc.Popover(
                [
                    dbc.PopoverBody([
                        dbc.Button("📊 Excel (.xlsx)", id="btn-export-excel",
                                   color="link", size="sm",
                                   style={"display": "block", "width": "100%",
                                          "textAlign": "left", "padding": "6px 0"}),
                        dbc.Button("📑 PowerPoint (.pptx)", id="btn-export-pptx",
                                   color="link", size="sm",
                                   style={"display": "block", "width": "100%",
                                          "textAlign": "left", "padding": "6px 0"}),
                        dcc.Download(id="download-excel"),
                        dcc.Download(id="download-pptx"),
                    ])
                ],
                target="btn-export-menu",
                trigger="click",
                placement="bottom-end",
            ),
        ],
    )


def _ai_chat():
    return html.Div([
        html.Button(
            "✨",
            id="fab-ai",
            className="ai-chat-fab",
            title="AI Assistant",
        ),
        dbc.Offcanvas(
            id="ai-offcanvas",
            title="AI Assistant",
            placement="end",
            is_open=False,
            style={"width": "400px"},
            children=[
                html.Div(
                    id="ai-messages",
                    className="ai-chat-messages",
                    style={"flex": 1, "overflowY": "auto", "padding": "12px 16px",
                           "display": "flex", "flexDirection": "column", "gap": "8px",
                           "minHeight": "300px"},
                    children=[
                        html.Div(
                            "Ciao! Sono il tuo assistente AI per l'analisi del processo. "
                            "Puoi chiedermi aiuto in qualsiasi momento — so sempre cosa "
                            "stai analizzando.",
                            className="chat-bubble-ai",
                        )
                    ],
                ),
                html.Div(
                    id="ai-chips-container",
                    className="chat-chips",
                ),
                html.Div(
                    [
                        dbc.Input(
                            id="ai-input",
                            placeholder="Scrivi un messaggio...",
                            type="text",
                            style={"fontSize": "13px", "flex": 1},
                            debounce=False,
                            n_submit=0,
                        ),
                        dbc.Button(
                            "↑",
                            id="btn-ai-send",
                            color="primary",
                            style={"borderRadius": "8px", "padding": "6px 14px",
                                   "fontWeight": 700},
                        ),
                    ],
                    style={"display": "flex", "gap": "8px", "padding": "12px 16px",
                           "borderTop": "1px solid #E2E8F0"},
                ),
                dcc.Store(id="ai-loading-flag", data=False),
            ],
        ),
    ])


app.layout = html.Div([
    # ── Session stores ──────────────────────────────────────────
    dcc.Store(id="ctx-store",   storage_type="session", data={}),
    dcc.Store(id="cfg-store",   storage_type="session",
              data={"alpha": 0.95, "y_cols": [], "excl_cols": []}),
    dcc.Store(id="data-store",  storage_type="memory",  data={}),
    dcc.Store(id="split-store", storage_type="memory",  data={}),
    dcc.Store(id="model-store", storage_type="memory",  data={}),
    dcc.Store(id="mon-store",   storage_type="memory",
              data={"mon_files": [], "anomaly_log": []}),
    dcc.Store(id="chat-store",  storage_type="session", data={"history": []}),
    dcc.Store(id="url-store",   storage_type="memory",  data="/"),

    # ── App Shell ───────────────────────────────────────────────
    dcc.Location(id="url", refresh=False),
    _header(),
    _sidebar(),
    html.Div(dash.page_container, id="main-content"),
    _ai_chat(),
])


# ── Nav active state ─────────────────────────────────────────
@callback(
    Output("url-store", "data"),
    Input("url", "pathname"),
)
def store_url(pathname):
    return pathname or "/"


# ── Header: process name + status badge ───────────────────────
@callback(
    Output("header-process-name", "children"),
    Output("header-status-badge", "children"),
    Input("ctx-store",   "data"),
    Input("model-store", "data"),
    Input("mon-store",   "data"),
)
def update_header(ctx, model, mon):
    ctx = ctx or {}
    model = model or {}
    mon = mon or {}

    # Process name snippet
    proc = ctx.get("process_context", "")
    name_chip = html.Span(
        f"— {proc[:40]}{'…' if len(proc) > 40 else ''}" if proc else "",
        style={"fontSize": "12px", "color": "#94A3B8"},
    ) if proc else ""

    # Status badge
    if not model.get("model_b64"):
        badge = html.Span("⚪ Nessun modello", className="badge badge-muted")
    else:
        mon_files = mon.get("mon_files", [])
        if mon_files:
            import numpy as np
            from core.store_utils import b64_to_model, mon_dict_to_mon
            mdl = b64_to_model(model["model_b64"])
            all_flags = []
            for mf in mon_files:
                m = mon_dict_to_mon(mf)
                all_flags.append(m["T2_flag"] | m["Q_flag"])
            flags = np.concatenate(all_flags)
            pct = flags.mean() * 100
            if pct < 5:
                badge = html.Span("🟢 In Control", className="badge badge-ok")
            elif pct < 15:
                badge = html.Span("🟡 Warning", className="badge badge-warn")
            else:
                badge = html.Span("🔴 Anomalie", className="badge badge-alarm")
        else:
            badge = html.Span("🔵 Modello calibrato", className="badge badge-info")

    return name_chip, badge


# ── Nav status icons ─────────────────────────────────────────
@callback(
    Output({"type": "nav-status", "page": "/"}, "children"),
    Output({"type": "nav-status", "page": "/setup"}, "children"),
    Output({"type": "nav-status", "page": "/training"}, "children"),
    Output({"type": "nav-status", "page": "/structure"}, "children"),
    Output({"type": "nav-status", "page": "/monitoring"}, "children"),
    Output({"type": "nav-status", "page": "/analysis"}, "children"),
    Input("ctx-store",   "data"),
    Input("data-store",  "data"),
    Input("model-store", "data"),
    Input("mon-store",   "data"),
)
def update_nav_status(ctx, data, model, mon):
    ctx   = ctx   or {}
    data  = data  or {}
    model = model or {}
    mon   = mon   or {}

    setup_ok    = bool(ctx.get("context_saved") and data.get("df_X_json"))
    model_ok    = bool(model.get("model_b64"))
    mon_ok      = bool(mon.get("mon_files"))

    return (
        "",                    # dashboard — always accessible
        "✅" if setup_ok else "⏳",
        "✅" if model_ok else ("⏳" if setup_ok else "🔒"),
        "⏳" if model_ok else "🔒",
        "⏳" if model_ok else "🔒",
        "✅" if mon_ok else ("⏳" if model_ok else "🔒"),
    )


# ── AI FAB → open offcanvas ───────────────────────────────────
@callback(
    Output("ai-offcanvas", "is_open"),
    Input("fab-ai",         "n_clicks"),
    State("ai-offcanvas",  "is_open"),
    prevent_initial_call=True,
)
def toggle_ai(n, is_open):
    return not is_open


# ── AI chips — contextual per page ───────────────────────────
@callback(
    Output("ai-chips-container", "children"),
    Input("url-store", "data"),
)
def update_chips(pathname):
    chips = PAGE_CHIPS.get(pathname or "/", [])
    return [
        html.Button(c, className="chip", id={"type": "ai-chip", "idx": i},
                    n_clicks=0)
        for i, c in enumerate(chips)
    ]


# ── AI send (button click or Enter) ──────────────────────────
@callback(
    Output("ai-messages",  "children"),
    Output("chat-store",   "data"),
    Output("ai-input",     "value"),
    Input("btn-ai-send",   "n_clicks"),
    Input("ai-input",      "n_submit"),
    Input({"type": "ai-chip", "idx": dash.ALL}, "n_clicks"),
    State("ai-input",      "value"),
    State("ai-messages",   "children"),
    State("chat-store",    "data"),
    State("ctx-store",     "data"),
    State("data-store",    "data"),
    State("model-store",   "data"),
    State("mon-store",     "data"),
    State("url-store",     "data"),
    prevent_initial_call=True,
)
def send_ai_message(n_btn, n_enter, chip_clicks, user_text,
                    current_messages, chat_data,
                    ctx, data, model, mon, pathname):
    from dash import ctx as dash_ctx
    import json

    # Determine the message text (button, enter, or chip)
    triggered = dash_ctx.triggered_id
    text = user_text or ""

    if isinstance(triggered, dict) and triggered.get("type") == "ai-chip":
        idx = triggered["idx"]
        chips = PAGE_CHIPS.get(pathname or "/", [])
        if idx < len(chips):
            text = chips[idx]

    if not text or not text.strip():
        return no_update, no_update, no_update

    text = text.strip()
    history = (chat_data or {}).get("history", [])

    # Build silent context snapshot
    ctx = ctx or {}
    data = data or {}
    model = model or {}
    mon = mon or {}
    lines = [f"[Pagina corrente]: {pathname or '/'}"]
    if ctx.get("process_context"):
        lines.append(f"[Processo]: {ctx['process_context'][:200]}")
    if ctx.get("analysis_objective"):
        lines.append(f"[Obiettivo]: {ctx['analysis_objective']}")
    if data.get("feature_names"):
        fn = data["feature_names"]
        lines.append(f"[Variabili X]: {', '.join(fn[:6])}{'...' if len(fn) > 6 else ''}")
    if model.get("model_b64"):
        from core.store_utils import b64_to_model
        try:
            mdl = b64_to_model(model["model_b64"])
            import numpy as np
            nf = int(((mdl['T2'] > mdl['T2_UCL']) | (mdl['Q'] > mdl['Q_UCL'])).sum())
            lines.append(
                f"[Modello]: k={mdl['k']} PC, T²_UCL={mdl['T2_UCL']:.3f}, "
                f"Q_UCL={mdl['Q_UCL']:.3f}, flag Phase I={nf}"
            )
        except Exception:
            pass
    if mon.get("mon_files"):
        import numpy as np
        from core.store_utils import mon_dict_to_mon
        mfs = mon["mon_files"]
        try:
            flags_all = np.concatenate([
                mon_dict_to_mon(mf)["T2_flag"] | mon_dict_to_mon(mf)["Q_flag"]
                for mf in mfs
            ])
            lines.append(
                f"[Monitoring]: {len(flags_all)} cicli, "
                f"{int(flags_all.sum())} anomalie ({flags_all.mean()*100:.1f}%)"
            )
        except Exception:
            pass

    snapshot = "\n".join(lines)
    lang = detect_language(ctx.get("process_context", "")) if ctx.get("process_context") else "Italian"

    full_prompt = f"{snapshot}\n\n---\nUtente: {text}"
    response_text, model_used, err = call_ai(full_prompt, language=lang)
    if err or not response_text:
        response_text = f"❌ Errore AI: {err or 'nessuna risposta'}"
        model_used = ""

    # Update history
    history = history + [
        {"role": "user",      "content": text},
        {"role": "assistant", "content": response_text},
    ]

    # Rebuild message list
    new_messages = list(current_messages or [])
    new_messages.append(html.Div(text, className="chat-bubble-user"))
    new_messages.append(
        html.Div([
            dcc.Markdown(response_text, style={"margin": 0, "fontSize": "13px"}),
            html.Div(
                f"via {model_used}" if model_used else "",
                style={"fontSize": "10px", "color": "#94A3B8", "marginTop": "4px"},
            ),
        ], className="chat-bubble-ai")
    )

    return new_messages, {"history": history[-40:]}, ""


# ── Session save ──────────────────────────────────────────────
@callback(
    Output("download-session", "data"),
    Input("btn-save-session",  "n_clicks"),
    State("ctx-store",   "data"),
    State("cfg-store",   "data"),
    State("data-store",  "data"),
    State("split-store", "data"),
    State("model-store", "data"),
    State("mon-store",   "data"),
    prevent_initial_call=True,
)
def save_session(n, ctx, cfg, data, split, model, mon):
    import pickle, base64, json
    bundle = dict(ctx=ctx, cfg=cfg, data=data, split=split, model=model, mon=mon)
    b64 = base64.b64encode(pickle.dumps(bundle)).decode()
    return dict(content=b64, filename="process_monitor_session.pkl",
                base64=True, type="application/octet-stream")


# ── Session load ──────────────────────────────────────────────
@callback(
    Output("ctx-store",        "data", allow_duplicate=True),
    Output("cfg-store",        "data", allow_duplicate=True),
    Output("data-store",       "data", allow_duplicate=True),
    Output("split-store",      "data", allow_duplicate=True),
    Output("model-store",      "data", allow_duplicate=True),
    Output("mon-store",        "data", allow_duplicate=True),
    Output("session-load-msg", "children"),
    Input("upload-session",    "contents"),
    State("upload-session",    "filename"),
    prevent_initial_call=True,
)
def load_session(contents, filename):
    if not contents:
        return (no_update,) * 6 + (no_update,)
    try:
        import pickle, base64
        _, content_string = contents.split(',')
        bundle = pickle.loads(base64.b64decode(content_string))
        msg = html.Span("✅ Sessione caricata",
                        style={"fontSize": "11px", "color": "#16A34A"})
        return (
            bundle.get("ctx",   {}),
            bundle.get("cfg",   {"alpha": 0.95, "y_cols": [], "excl_cols": []}),
            bundle.get("data",  {}),
            bundle.get("split", {}),
            bundle.get("model", {}),
            bundle.get("mon",   {"mon_files": [], "anomaly_log": []}),
            msg,
        )
    except Exception as e:
        msg = html.Span(f"❌ Errore: {e}", style={"fontSize": "11px", "color": "#DC2626"})
        return (no_update,) * 6 + (msg,)


# ── Excel export ──────────────────────────────────────────────
@callback(
    Output("download-excel",    "data"),
    Input("btn-export-excel",   "n_clicks"),
    State("data-store",  "data"),
    State("split-store", "data"),
    State("model-store", "data"),
    State("mon-store",   "data"),
    prevent_initial_call=True,
)
def export_excel(n, data, split, model_data, mon):
    if not model_data or not model_data.get("model_b64"):
        return no_update
    try:
        import io as _io
        import numpy as np
        import pandas as pd
        from core.store_utils import b64_to_model, json_to_df, mon_dict_to_mon

        mdl = b64_to_model(model_data["model_b64"])
        fn = mdl["feature_names"]
        buf = _io.BytesIO()

        with pd.ExcelWriter(buf, engine="openpyxl") as writer:
            # Model summary
            evr_cum = np.cumsum(mdl["evr"])
            pd.DataFrame({
                "Metric": ["k PCs", "T² UCL", "Q UCL", "N calibration",
                           f"Variance PC1–{mdl['k']} (%)"],
                "Value":  [mdl["k"], round(mdl["T2_UCL"], 4), round(mdl["Q_UCL"], 4),
                           mdl["n_train"], round(float(evr_cum[mdl["k"] - 1]), 2)],
            }).to_excel(writer, sheet_name="Model_Summary", index=False)

            # Calibration data
            if data.get("df_X_json"):
                df_X = json_to_df(data["df_X_json"])
                cal_use = json_to_df(split["df_train_json"]) if split.get("df_train_json") else df_X
                df_cal = cal_use.copy()
                df_cal["T²"] = mdl["T2"].round(4)
                df_cal["Q"]  = mdl["Q"].round(4)
                df_cal["T²_flag"] = mdl["T2"] > mdl["T2_UCL"]
                df_cal["Q_flag"]  = mdl["Q"]  > mdl["Q_UCL"]
                df_cal.to_excel(writer, sheet_name="Calibration", index=False)

            # Monitoring data
            mfiles = (mon or {}).get("mon_files", [])
            if mfiles:
                rows = []
                for mf in mfiles:
                    m = mon_dict_to_mon(mf)
                    for i in range(mf["n_rows"]):
                        rows.append({
                            "File": mf["name"], "Cycle": i,
                            "T²": round(float(m["T2"][i]), 4),
                            "Q":  round(float(m["Q"][i]),  4),
                            "T²_flag": bool(m["T2_flag"][i]),
                            "Q_flag":  bool(m["Q_flag"][i]),
                            "Severity×UCL": round(float(
                                max(m["T2"][i] / mdl["T2_UCL"],
                                    m["Q"][i]  / mdl["Q_UCL"])
                            ), 2),
                        })
                pd.DataFrame(rows).to_excel(writer, sheet_name="Monitoring", index=False)

                # Anomalies
                all_t2 = np.concatenate([mon_dict_to_mon(mf)["T2"] for mf in mfiles])
                all_q  = np.concatenate([mon_dict_to_mon(mf)["Q"]  for mf in mfiles])
                flagged = np.where(
                    (all_t2 > mdl["T2_UCL"]) | (all_q > mdl["Q_UCL"])
                )[0]
                if len(flagged):
                    pd.DataFrame({
                        "Cycle": flagged,
                        "T²":    all_t2[flagged].round(4),
                        "T²/UCL": (all_t2[flagged] / mdl["T2_UCL"]).round(3),
                        "Q":     all_q[flagged].round(4),
                        "Q/UCL": (all_q[flagged] / mdl["Q_UCL"]).round(3),
                        "Severity×UCL": np.maximum(
                            all_t2[flagged] / mdl["T2_UCL"],
                            all_q[flagged]  / mdl["Q_UCL"],
                        ).round(3),
                    }).sort_values("Severity×UCL", ascending=False
                    ).to_excel(writer, sheet_name="Anomalies", index=False)

            # Loadings
            P = mdl["loadings"]
            df_load = pd.DataFrame(
                P, columns=[f"PC{i + 1}" for i in range(P.shape[1])]
            )
            df_load.insert(0, "Variable", fn)
            df_load.to_excel(writer, sheet_name="Loadings", index=False)

            # Anomaly log
            alog = (mon or {}).get("anomaly_log", [])
            if alog:
                pd.DataFrame(alog).to_excel(writer, sheet_name="Intervention_Log", index=False)

        buf.seek(0)
        import base64 as _b64
        content = _b64.b64encode(buf.read()).decode()
        return dict(content=content, filename="process_monitor_export.xlsx",
                    base64=True,
                    type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    except Exception as e:
        print(f"Excel export error: {e}")
        return no_update


# ── PowerPoint export ─────────────────────────────────────────
@callback(
    Output("download-pptx",   "data"),
    Input("btn-export-pptx",  "n_clicks"),
    State("ctx-store",   "data"),
    State("model-store", "data"),
    State("mon-store",   "data"),
    prevent_initial_call=True,
)
def export_pptx(n, ctx_data, model_data, mon):
    if not model_data or not model_data.get("model_b64"):
        return no_update
    try:
        import io as _io
        import base64 as _b64
        import numpy as np
        from datetime import date
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from core.store_utils import b64_to_model, mon_dict_to_mon
        from components.charts import chart_line_multi, COLORS as CC
        import plotly.io as pio

        mdl = b64_to_model(model_data["model_b64"])
        fn = mdl["feature_names"]
        ctx_data = ctx_data or {}
        mfiles = (mon or {}).get("mon_files", [])

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        def add_slide(title_text, subtitle_text=""):
            slide = prs.slides.add_slide(blank_layout)
            tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
            tf.text_frame.text = title_text
            tf.text_frame.paragraphs[0].runs[0].font.size = Pt(22)
            tf.text_frame.paragraphs[0].runs[0].font.bold = True
            tf.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            if subtitle_text:
                tf2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12), Inches(0.4))
                tf2.text_frame.text = subtitle_text
                tf2.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
                tf2.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
            return slide

        def img_from_fig(fig, w=900, h=350):
            img_bytes = pio.to_image(fig, format="png", width=w, height=h, scale=2)
            return _io.BytesIO(img_bytes)

        # Slide 1: Title
        s1 = add_slide("Process Monitor — Report Analisi")
        proc = ctx_data.get("process_context", "Processo industriale")
        obj  = ctx_data.get("analysis_objective", "")
        body_tf = s1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(4))
        body_tf.text_frame.word_wrap = True
        for line in [
            f"Processo: {proc[:120]}",
            f"Obiettivo: {obj[:100]}" if obj else "",
            f"Data: {date.today().isoformat()}",
            f"Modello: k={mdl['k']} PC, α=0.95",
        ]:
            if line:
                p = body_tf.text_frame.add_paragraph()
                p.text = line
                p.runs[0].font.size = Pt(14)

        # Slide 2: Model KPIs
        evr_cum = float(np.cumsum(mdl["evr"])[mdl["k"] - 1])
        nf_p1 = int(((mdl["T2"] > mdl["T2_UCL"]) | (mdl["Q"] > mdl["Q_UCL"])).sum())
        s2 = add_slide("Modello PCA-SPC — Parametri chiave")
        kpis = [
            ("Componenti k", str(mdl["k"])),
            ("UCL T²", f"{mdl['T2_UCL']:.3f}"),
            ("UCL Q (SPE)", f"{mdl['Q_UCL']:.3f}"),
            ("Varianza spiegata", f"{evr_cum:.1f}%"),
            ("Cicli calibrazione", str(mdl["n_train"])),
            ("Flag Phase I", f"{nf_p1} ({nf_p1/len(mdl['T2'])*100:.1f}%)"),
        ]
        for i, (lbl, val) in enumerate(kpis):
            col = i % 3
            row = i // 3
            x = Inches(0.5 + col * 4.2)
            y = Inches(1.5 + row * 1.8)
            tb = s2.shapes.add_textbox(x, y, Inches(3.8), Inches(1.5))
            tb.text_frame.text = val
            tb.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
            tb.text_frame.paragraphs[0].runs[0].font.bold = True
            tb.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            p2 = tb.text_frame.add_paragraph()
            p2.text = lbl
            p2.runs[0].font.size = Pt(11)
            p2.runs[0].font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

        # Slides 3-4: Control charts (only if monitoring data exists)
        if mfiles:
            all_t2 = np.concatenate([mon_dict_to_mon(mf)["T2"]      for mf in mfiles])
            all_q  = np.concatenate([mon_dict_to_mon(mf)["Q"]       for mf in mfiles])
            all_t2f= np.concatenate([mon_dict_to_mon(mf)["T2_flag"] for mf in mfiles])
            all_qf = np.concatenate([mon_dict_to_mon(mf)["Q_flag"]  for mf in mfiles])
            boundaries = [0] + list(np.cumsum([mf["n_rows"] for mf in mfiles]))
            names = [mf["name"] for mf in mfiles]

            for stat_arr, ucl_val, flag_arr, stat_name, color in [
                (all_t2, mdl["T2_UCL"], all_t2f, "Hotelling T²",  CC["primary"]),
                (all_q,  mdl["Q_UCL"],  all_qf,  "Q (SPE)",       CC["success"]),
            ]:
                fig = chart_line_multi(stat_arr, ucl_val, f"Phase II — {stat_name}",
                                       color, flag_arr, names, boundaries)
                fig.update_layout(height=350, width=900)
                s = add_slide(f"Control Chart — {stat_name}")
                img_stream = img_from_fig(fig)
                s.shapes.add_picture(img_stream, Inches(0.5), Inches(1.3),
                                     Inches(12.3), Inches(4.5))

            # Slide 5: Top anomalies table
            flagged_idx = np.where(all_t2f | all_qf)[0]
            if len(flagged_idx):
                n_any = len(flagged_idx)
                pct   = n_any / len(all_t2) * 100
                sev   = np.maximum(all_t2[flagged_idx] / mdl["T2_UCL"],
                                   all_q[flagged_idx]  / mdl["Q_UCL"])
                top5  = np.argsort(sev)[::-1][:5]
                stato = "STABILE" if pct < 5 else "WARNING" if pct < 15 else "ANOMALIE"
                s5 = add_slide(
                    f"Top anomalie per severity — {stato}",
                    f"{n_any} anomalie su {len(all_t2)} cicli ({pct:.1f}%)",
                )
                hdr = ["Ciclo", "T²", "T²/UCL", "Q", "Q/UCL", "Severity"]
                rows_data = []
                for idx_ in flagged_idx[top5]:
                    rows_data.append([
                        str(int(idx_)),
                        f"{all_t2[idx_]:.3f}",
                        f"{all_t2[idx_]/mdl['T2_UCL']:.2f}×",
                        f"{all_q[idx_]:.3f}",
                        f"{all_q[idx_]/mdl['Q_UCL']:.2f}×",
                        f"{max(all_t2[idx_]/mdl['T2_UCL'], all_q[idx_]/mdl['Q_UCL']):.2f}×",
                    ])
                col_widths = [Inches(1.5), Inches(1.8), Inches(1.8),
                              Inches(1.8), Inches(1.8), Inches(2.0)]
                rows_count = len(rows_data) + 1
                tbl = s5.shapes.add_table(
                    rows_count, len(hdr), Inches(0.5), Inches(1.5),
                    Inches(12.3), Inches(0.45 * rows_count)
                ).table
                for ci, h in enumerate(hdr):
                    cell = tbl.cell(0, ci)
                    cell.text = h
                    cell.text_frame.paragraphs[0].runs[0].font.bold = True
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                for ri, row in enumerate(rows_data, start=1):
                    for ci, val in enumerate(row):
                        cell = tbl.cell(ri, ci)
                        cell.text = val
                        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

        # Save
        pptx_buf = _io.BytesIO()
        prs.save(pptx_buf)
        pptx_buf.seek(0)
        content = _b64.b64encode(pptx_buf.read()).decode()
        return dict(
            content=content,
            filename="process_monitor_report.pptx",
            base64=True,
            type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as e:
        print(f"PPTX export error: {e}")
        return no_update


if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=8050)
